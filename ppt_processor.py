try:
    # 导入 LibreOffice 导出函数
    from ppt_exporter_libreoffice import export_slides_with_libreoffice
    LIBREOFFICE_EXPORTER_AVAILABLE = True
except ImportError:
    logging.warning("无法导入 'ppt_exporter_libreoffice.py' 或其依赖。LibreOffice 导出功能将不可用。")
    LIBREOFFICE_EXPORTER_AVAILABLE = False
    def export_slides_with_libreoffice(pptx_filepath, output_dir): # 定义空函数
        logging.error("尝试调用 LibreOffice 导出功能，但模块未加载。")
        return None

import os
import platform
import time
import shutil
# import wave
# import contextlib
from pathlib import Path
import logging
import uuid
import configparser # 导入配置解析器
import json # 需要 json 模块解析 ffprobe 输出
import subprocess # 需要 subprocess 调用 ffprobe
import shlex # 需要 shlex 处理命令行参数
# 确保导入 mutagen
try:
    from mutagen.mp3 import MP3
    from mutagen import File as MutagenFile, MutagenError
    MUTAGEN_AVAILABLE = True
except ImportError:
    logging.warning("缺少 'mutagen' 库，MP3 时长可能不准。'pip install mutagen'")
    MUTAGEN_AVAILABLE = False


# !!! CHANGE: Import the new TTS manager !!!
try:
    import tts_manager_edge as tts_manager # 使用别名
except ImportError:
     logging.error("无法导入 tts_manager_edge.py!")
     exit()

try:
    from ppt_exporter_win import export_slides_with_powerpoint
    WINDOWS_EXPORTER_AVAILABLE = True
except ImportError:
    # 即使在非 Windows 平台，这个 warning 也没关系，只是表示该功能不可用
    logging.debug("无法导入 'ppt_exporter_win.py'。Windows 平台导出功能将不可用。")
    WINDOWS_EXPORTER_AVAILABLE = False
    def export_slides_with_powerpoint(pptx_filepath, output_dir): # 定义空函数
        # logging.error("尝试调用 Windows 导出功能，但模块未加载。") # 避免在非 Windows 下报错
        return None


# --- 配置解析 ---
config = configparser.ConfigParser()
# 读取配置文件，如果文件不存在或读取失败，将使用下面定义的默认值
# 使用 'utf-8' 编码以支持中文路径或值
config = configparser.ConfigParser()
config_path = Path(__file__).parent / 'config.ini'
if config_path.exists():
    try: config.read(config_path, encoding='utf-8'); logging.info(f"加载配置: {config_path}")
    except Exception as e: logging.error(f"加载配置失败: {e}")
else: logging.warning(f"配置未找到: {config_path}")


# --- 日志记录配置 (现在可以从 config 读取级别) ---
log_level_str = config.get('General', 'logging_level', fallback='INFO').upper()
log_level = getattr(logging, log_level_str, logging.INFO) # 获取对应级别，无效则默认为 INFO
logging.basicConfig(
    level=log_level,
    format='%(asctime)s - %(levelname)s - [%(filename)s:%(lineno)d] - %(message)s',
    datefmt='%Y-%m-%d %H:%M:%S'
)


# --- 导入其他模块 (放在日志配置之后) ---
try:
    from ppt_exporter_win import export_slides_with_powerpoint
except ImportError:
    logging.warning("无法导入 'ppt_exporter_win.py'。Windows 平台导出功能将不可用。")
    def export_slides_with_powerpoint(pptx_filepath, output_dir): # 定义一个空函数避免 NameError
        logging.error("尝试调用 Windows 导出功能，但模块未加载。")
        return None
try:
    from pptx import Presentation
except ImportError:
    logging.error("缺少 'python-pptx' 库。请使用 'pip install python-pptx' 安装。")
    exit()
try:
    import pyttsx3
except ImportError:
    logging.error("缺少 'pyttsx3' 库。请使用 'pip install pyttsx3' 安装。")
    exit()


# --- TTS 相关配置 (从 config 读取) ---
# TTS_RATE = config.getint('Audio', 'tts_rate', fallback=180)
TTS_RATE_PERCENT = config.getint('Audio', 'tts_rate_percent', fallback=100) # 从配置读速率百分比
# TTS_VOICE_ID = config.get('Audio', 'tts_voice_id', fallback=None)
# # 如果读取到的是空字符串，也视为 None
# if not TTS_VOICE_ID:
#     TTS_VOICE_ID = None


# --- 获取 FFmpeg/FFprobe 路径 (需要确保能找到 ffprobe) ---
def get_ffmpeg_tool_path(tool_name="ffmpeg"):
    """确定 ffmpeg 或 ffprobe 的路径"""
    # 优先从 config.ini 读取 tool_name 对应的路径
    tool_path_config = config.get('Paths', f'{tool_name}_path', fallback=tool_name)

    # 尝试在 PATH 或配置路径中查找
    tool_executable_found = shutil.which(tool_path_config)
    if tool_executable_found:
        logging.info(f"找到 {tool_name} 可执行文件: {tool_executable_found}")
        return str(Path(tool_executable_found).resolve())

    # 如果找不到，尝试在 ffmpeg 所在的目录下查找
    ffmpeg_path_str = config.get('Paths', 'ffmpeg_path', fallback='ffmpeg')
    ffmpeg_found = shutil.which(ffmpeg_path_str)
    if ffmpeg_found:
        ffmpeg_dir = Path(ffmpeg_found).parent
        tool_in_ffmpeg_dir = ffmpeg_dir / f"{tool_name}.exe" # 假设是 windows
        if tool_in_ffmpeg_dir.exists():
             logging.info(f"在 FFmpeg 目录中找到 {tool_name}: {tool_in_ffmpeg_dir}")
             return str(tool_in_ffmpeg_dir.resolve())

    logging.error(f"未能找到 {tool_name} 可执行文件！请确保已安装或在 config.ini 中配置路径。")
    return None

FFMPEG_PATH_RESOLVED = get_ffmpeg_tool_path("ffmpeg")
FFPROBE_PATH_RESOLVED = get_ffmpeg_tool_path("ffprobe")


# --- 工具函数 (get_wav_duration - 保持不变) ---
## --- 工具函数 (get_audio_duration - 使用 FFprobe 重写) ---
def get_audio_duration(filepath: Path) -> float | None:
    """
    使用 FFprobe 获取音频文件的准确时长 (秒)。

    Args:
        filepath: 音频文件的 Path 对象。

    Returns:
        音频时长 (float)，如果无法获取则返回 None。
    """
    if not filepath or not filepath.is_file():
        logging.warning(f"尝试获取时长失败，文件无效或不存在: {filepath}")
        return None
    if FFPROBE_PATH_RESOLVED is None:
        logging.error("无法获取音频时长，因为找不到 ffprobe。")
        return None

    command = [
        FFPROBE_PATH_RESOLVED,
        "-v", "quiet", # 静默模式，只输出错误或指定信息
        "-print_format", "json", # 输出为 JSON 格式
        "-show_format", # 显示格式信息，包含时长
        "-show_streams", # 显示流信息 (有时时长在流信息里)
        str(filepath.resolve())
    ]

    try:
        logging.debug(f"执行 ffprobe 获取时长: {shlex.join(command)}")
        result = subprocess.run(command, capture_output=True, text=True, check=True, encoding='utf-8', errors='ignore')
        metadata = json.loads(result.stdout) # 解析 JSON 输出

        duration = None
        # 优先从 format -> duration 获取
        if 'format' in metadata and 'duration' in metadata['format']:
            try:
                duration = float(metadata['format']['duration'])
                logging.debug(f"从 format 获取 {filepath.name} 时长: {duration:.3f}s")
            except (ValueError, TypeError):
                 logging.warning(f"无法从 format.duration 解析 {filepath.name} 的有效时长: {metadata['format'].get('duration')}")

        # 如果 format 中没有，尝试从第一个音频流的 duration 获取
        if duration is None and 'streams' in metadata:
            for stream in metadata['streams']:
                if stream.get('codec_type') == 'audio' and 'duration' in stream:
                    try:
                        duration = float(stream['duration'])
                        logging.debug(f"从 audio stream 获取 {filepath.name} 时长: {duration:.3f}s")
                        break # 找到第一个音频流的时长就够了
                    except (ValueError, TypeError):
                        logging.warning(f"无法从 stream.duration 解析 {filepath.name} 的有效时长: {stream.get('duration')}")

        if duration is not None and duration >= 0: # 时长不能是负数
             # 对极短时长进行判断
             if duration < 0.01:
                  logging.warning(f"FFprobe 获取的时长过短 ({duration:.3f}s) for {filepath.name}，可能无效。")
                  # 可以返回 0.0 或 None，返回 None 让调用者决定如何处理
                  return None
             return duration
        else:
            logging.error(f"FFprobe未能从 {filepath.name} 的元数据中找到有效的时长信息。元数据: {metadata}")
            return None

    except subprocess.CalledProcessError as e:
        logging.error(f"执行 ffprobe 失败 for {filepath.name}。返回码: {e.returncode}")
        logging.error(f"FFprobe 命令: {shlex.join(command)}")
        logging.error(f"FFprobe 错误输出:\n{e.stderr}")
        return None
    except json.JSONDecodeError as e:
        logging.error(f"解析 ffprobe 的 JSON 输出失败 for {filepath.name}: {e}")
        logging.error(f"FFprobe 原始输出:\n{result.stdout}")
        return None
    except FileNotFoundError:
        logging.error(f"错误：找不到 ffprobe 命令 '{FFPROBE_PATH_RESOLVED}'。")
        return None
    except Exception as e:
        logging.error(f"使用 ffprobe 获取 {filepath.name} 时长时发生未知错误: {e}", exc_info=True)
        return None

# --- 核心处理函数 ---

def extract_speaker_notes(pptx_filepath: Path) -> list[str] | None:
    # ... (函数体保持不变) ...
    if not pptx_filepath.is_file():
        logging.error(f"输入文件不存在: {pptx_filepath}")
        return None

    notes_list = []
    try:
        logging.info(f"开始解析演示文稿以提取备注: {pptx_filepath.name}")
        prs = Presentation(pptx_filepath)
        num_slides = len(prs.slides)
        logging.info(f"演示文稿包含 {num_slides} 张幻灯片。")

        for i, slide in enumerate(prs.slides):
            slide_num = i + 1
            note_text = "" # 默认为空字符串
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                if text_frame and text_frame.text:
                    note_text = text_frame.text.strip()
                    logging.debug(f"  找到幻灯片 {slide_num} 的备注: '{note_text[:50]}...'") # 记录备注开头部分
            else:
                logging.debug(f"  幻灯片 {slide_num} 没有备注。")
            notes_list.append(note_text)

        logging.info(f"成功提取了 {len(notes_list)} 条备注信息。")
        return notes_list

    except Exception as e:
        logging.error(f"解析 PPTX 文件以提取备注时出错: {e}", exc_info=True) # exc_info=True 会记录堆栈跟踪
        return None


# --- generate_audio_segments 函数 (修改调用) ---
def generate_audio_segments(
    notes: list[str],
    output_audio_dir: Path,
    voice_id: str # voice_id 现在是必需的
# ) -> list[tuple[str | None, float | None]]: # 返回类型不变 (仍然可能返回 None 时长)
) -> list[tuple[str | None, float | None]]:
    """
    使用 Edge TTS 将文本备注列表转换为 MP3 音频文件。

    Args:
        notes: 包含每张幻灯片备注文本的字符串列表。
        output_audio_dir: 保存生成的 MP3 文件的目标目录。
        voice_id: 要使用的 Edge TTS 语音 ID (必需)。

    Returns:
        一个元组列表，每个元组包含 (生成的音频文件路径 | None, 音频时长)。
    """
    audio_results = []
    output_audio_dir.mkdir(parents=True, exist_ok=True)

    # 从配置获取速率 (全局)
    rate_percent = TTS_RATE_PERCENT
    logging.info(f"将使用 Edge TTS 速率: {rate_percent}% (来自配置)")

    logging.info(f"开始使用 Edge TTS 生成音频片段 (Voice ID: {voice_id})...")
    total_duration = 0.0
    audio_results = [] # 移到循环外初始化

    for i, text in enumerate(notes):
        segment_num = i + 1
        # !!! CHANGE: Output format is now MP3 !!!
        audio_filename = f"segment_{segment_num}.mp3"
        audio_filepath = output_audio_dir / audio_filename
        # !!! 修改: 初始化 duration_sec 为 None !!!
        duration_sec = None
        result_path = None

        if not text or text.isspace():
            logging.info(f"  片段 {segment_num}: 文本为空，跳过 TTS。")
            # !!! 修改: 时长也记录为 None !!!
            audio_results.append((None, None))
            continue

        logging.info(f"  正在生成片段 {segment_num} 的音频 (文本: '{text[:30]}...')...")
        # !!! CHANGE: Call generate_segment_audio without pitch !!!
        success = tts_manager.generate_segment_audio(
            voice_id,
            text,
            audio_filepath,
            rate=rate_percent
        )

        # !!! --------------------------------------------- !!!

        if success:
            # --- 关键: 调用新的 get_audio_duration ---
            duration_sec = get_audio_duration(audio_filepath) # 可能返回 None 或 float
            # --- ----------------------------------- ---
            # 获取时长
            # duration_sec = get_audio_duration(audio_filepath) # duration_sec 可能为 None
            if duration_sec is not None: # <<< 关键检查：确保 duration_sec 不是 None
                # !!! 修改: 现在可以安全地进行比较了 !!!
                if duration_sec > 0.01:
                    result_path = str(audio_filepath.resolve())
                    logging.info(f"    片段 {segment_num} 音频已保存: {audio_filename} (时长: {duration_sec:.3f}s)")
                    total_duration += duration_sec # 只有有效时长才累加
                else: # 时长为 0 或过小
                    logging.warning(f"    片段 {segment_num} 文件已生成但有效时长为 0 或过短 ({duration_sec:.3f}s)。")
                    result_path = str(audio_filepath.resolve()) # 文件存在
                    duration_sec = 0.0 # 将其规范化为 0.0 用于后续处理
            else: # get_audio_duration 返回了 None
                logging.error(f"    无法获取片段 {segment_num} ({audio_filename}) 的有效时长!")
                result_path = str(audio_filepath.resolve()) # 文件可能存在，但时长未知
                # duration_sec 保持为 None
        else: # TTS 生成失败
            result_path = None
            duration_sec = None # TTS 失败，时长也为 None


        # !!! 修改: audio_results 中记录的时长可能是 None 或 float !!!
        audio_results.append((result_path, duration_sec))

    logging.info(f"所有音频片段生成完成。总预估旁白时长: {total_duration:.2f} 秒。")
    return audio_results


# --- process_presentation 函数 (修改后) ---
def process_presentation(
    pptx_filepath: Path,
    base_output_dir: Path,
    voice_id: str | None = None # 确保 voice_id 参数存在
) -> tuple[list[dict] | None, Path | None]:
    """
    完整的处理流程：导出幻灯片 -> 提取备注 -> 生成音频 (使用 Edge TTS)。
    """
    if not pptx_filepath.is_file():
        logging.error(f"输入 PPTX 文件不存在: {pptx_filepath}")
        return None, None

    if voice_id is None:
         logging.error("错误：调用 process_presentation 时未提供必需的 voice_id 参数！")
         # 或者可以尝试从全局配置回退，但不推荐
         # voice_id = TTS_VOICE_ID
         # if not voice_id:
         #    return None, None
         return None, None # 强制要求提供 voice_id

    run_id = uuid.uuid4().hex[:8]
    temp_run_dir = base_output_dir / f"temp_{pptx_filepath.stem}_{run_id}"
    temp_image_dir = temp_run_dir / "images"
    temp_audio_dir = temp_run_dir / "audio"

    try:
        temp_run_dir.mkdir(parents=True, exist_ok=True)
        logging.info(f"创建临时工作目录: {temp_run_dir}")
    except OSError as e:
         logging.error(f"无法创建临时工作目录 {temp_run_dir}: {e}")
         return None, None

    # --- 步骤 1: 导出幻灯片图片 ---
    logging.info("--- 步骤 1: 导出幻灯片图片 ---")
    image_paths = None
    exporter_used = "None" # 记录使用了哪个导出器

    if platform.system() == "Windows" and WINDOWS_EXPORTER_AVAILABLE:
        logging.info("检测到 Windows 平台，优先尝试使用 PowerPoint COM 接口导出...")
        try:
            image_paths = export_slides_with_powerpoint(pptx_filepath, temp_image_dir)
            if image_paths:
                logging.info("成功使用 PowerPoint COM 接口导出图片。")
                exporter_used = "PowerPoint COM"
            else:
                logging.warning("PowerPoint COM 导出失败或未返回路径。")
        except NameError: # 理论上 WINDOWS_EXPORTER_AVAILABLE 会处理，但双重保险
            logging.warning("Windows 导出函数未加载。")
        except Exception as e_com:
            logging.error(f"PowerPoint COM 导出时发生错误: {e_com}。")

        # 如果 PowerPoint 导出失败或不可用，尝试 LibreOffice 作为后备
        if not image_paths and LIBREOFFICE_EXPORTER_AVAILABLE:
            logging.warning("尝试使用 LibreOffice 作为后备导出方案...")
            try:
                image_paths = export_slides_with_libreoffice(pptx_filepath, temp_image_dir)
                if image_paths:
                    logging.info("成功使用 LibreOffice 导出图片。")
                    exporter_used = "LibreOffice"
                else:
                    logging.warning("LibreOffice 后备导出也失败了。")
            except Exception as e_lo_fallback:
                logging.error(f"LibreOffice 后备导出时发生错误: {e_lo_fallback}")

    # 如果不是 Windows，或者 Windows 导出器不可用，直接尝试 LibreOffice
    elif LIBREOFFICE_EXPORTER_AVAILABLE:
        logging.info("非 Windows 平台或 Windows 导出不可用，尝试使用 LibreOffice 导出...")
        try:
            image_paths = export_slides_with_libreoffice(pptx_filepath, temp_image_dir)
            if image_paths:
                logging.info("成功使用 LibreOffice 导出图片。")
                exporter_used = "LibreOffice"
            else:
                logging.warning("LibreOffice 导出失败。")
        except Exception as e_lo:
            logging.error(f"LibreOffice 导出时发生错误: {e_lo}")

    else:
        logging.error("没有可用的幻灯片导出方法 (PowerPoint COM 或 LibreOffice)。")

    # 检查最终结果
    if not image_paths:
        logging.error("所有幻灯片导出方法均失败。无法继续处理。")
        # 清理可能已创建的临时目录的一部分
        # if temp_run_dir.exists(): shutil.rmtree(temp_run_dir) # 这里清理可能过早
        return None, temp_run_dir # 返回失败，并传递临时目录路径供上层清理
    else:
        logging.info(f"使用 '{exporter_used}' 成功导出 {len(image_paths)} 张图片。")
    # --- 导出步骤修改结束 ---


    # --- 步骤 2: 提取备注 (保持不变) ---
    logging.info("--- 步骤 2: 提取演讲者备注 ---")
    notes_list = extract_speaker_notes(pptx_filepath)
    if notes_list is None:
        logging.error("提取备注失败。")
        return None, temp_run_dir

    # --- 步骤 3: 对齐图片和备注 (保持不变) ---
    # ... (对齐逻辑) ...
    num_images = len(image_paths)
    num_notes = len(notes_list)
    if num_images != num_notes:
        logging.warning(f"图片数({num_images})与备注数({num_notes})不匹配，将按较小数处理。")
        min_count = min(num_images, num_notes)
        image_paths = image_paths[:min_count]
        notes_list = notes_list[:min_count]


    # --- 步骤 4: 生成音频片段 ---
    logging.info("--- 步骤 4: 生成音频片段 (Edge TTS) ---")
    audio_results = generate_audio_segments(notes_list, temp_audio_dir, voice_id=voice_id)
    if len(audio_results) != len(notes_list):
         logging.error("TTS 结果数量不匹配！")
         return None, temp_run_dir
    logging.info("音频片段生成完成。")

    # --- 步骤 5: 组合结果 ---
    logging.info("--- 步骤 5: 整理处理结果 ---")
    final_data = []
    processing_successful = True

    for i in range(len(notes_list)):
        audio_path_str = audio_results[i][0] if i < len(audio_results) else None
        # !!! 修改: audio_duration 现在可能是 None 或 float !!!
        audio_duration_raw = audio_results[i][1] if i < len(audio_results) else None
        image_path_str = image_paths[i] if i < len(image_paths) else None

        # --- 处理时长 ---
        final_audio_duration = 0.0 # 默认最终使用的时长为 0
        if audio_duration_raw is not None: # 如果获取到了时长（包括 0）
            if audio_duration_raw > 0.01:
                final_audio_duration = audio_duration_raw # 使用有效时长
            else:
                 # 时长为 0 或过小，仍记录为 0
                 final_audio_duration = 0.0
                 if audio_path_str: # 仅当文件存在时记录警告
                     logging.warning(f"幻灯片 {i+1} 音频 ({Path(audio_path_str).name}) 时长为 0 或过短。")
        else: # audio_duration_raw is None (获取时长失败)
             if audio_path_str: # 文件存在但无法获取时长
                 logging.error(f"无法确定幻灯片 {i+1} 音频 ({Path(audio_path_str).name}) 的时长！将使用默认图片时长。")
                 # 标记处理可能不完全成功，但不中断
                 # processing_successful = False # 如果希望严格失败可以取消注释
             else: # 文件不存在 (TTS 失败或无备注)
                 pass # 时长为 0 是正常的
             # 最终时长仍为 0.0

        slide_data = {
            'slide_number': i + 1,
            'image_path': image_path_str,
            'notes': notes_list[i],
            'audio_path': audio_path_str,
            # !!! 修改: 存储最终确定的 float 时长 (0.0 或有效时长) !!!
            'audio_duration': final_audio_duration
        }
        final_data.append(slide_data)
        logging.debug(f"  整理幻灯片 {i+1} 数据: Image={Path(image_path_str).name if image_path_str else 'N/A'}, Audio={Path(audio_path_str).name if audio_path_str else 'N/A'}, Duration={final_audio_duration:.3f}s (Raw: {audio_duration_raw})")

    # if not processing_successful: # 如果标记了失败
    #     logging.error("处理过程中未能获取部分音频文件的准确时长...")
        # return None, temp_run_dir

    logging.info(f"成功整理了 {len(final_data)} 张幻灯片的数据。")
    return final_data, temp_run_dir


# --- 示例用法 (主要用于单独测试此模块) ---
if __name__ == "__main__":
    # ... (测试逻辑不变) ...
    # 注意：现在测试需要网络连接
    import inspect # 确保导入 inspect
    logging.info("--- 测试 PPT Processor 模块 (使用 Edge TTS, 需要网络) ---")
    test_ppt_file = Path("智能短信分类平台方案.pptx") # <--- 修改为你测试用的 PPTX 文件
    test_output_dir = Path("./processor_edge_test_output") # 改个名字

    if not test_ppt_file.exists():
         print(f"错误：请将第 {inspect.currentframe().f_lineno - 4} 行的文件路径替换为实际 PPTX 文件。")
    else:
        test_output_dir.mkdir(parents=True, exist_ok=True)
        for item in test_output_dir.glob("temp_*"):
             if item.is_dir():
                 print(f"清理旧的测试目录: {item}")
                 shutil.rmtree(item)

        # 在测试时需要提供一个 voice_id
        test_voice_id = "zh-CN-XiaoxiaoNeural" # 或者从 KNOWN_EDGE_VOICES 选一个
        print(f"测试将使用 Voice ID: {test_voice_id}")

        result_data, result_temp_dir = process_presentation(
            test_ppt_file,
            test_output_dir,
            voice_id=test_voice_id # <<< 提供 voice_id
        )

        if result_data and result_temp_dir:
            print("\n--- 处理成功 ---")
            print(f"临时文件位于: {result_temp_dir.resolve()}")
            print(f"处理了 {len(result_data)} 张幻灯片的数据。")
        else:
            print("\n--- 处理失败 ---")
            if result_temp_dir:
                 print(f"临时文件可能保留在: {result_temp_dir.resolve()}")

    logging.info("--- PPT Processor 模块测试结束 ---")